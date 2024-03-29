import streamlit as st
import numpy as np
import pandas as pd
import altair as alt
import cv2
import easyocr
import os
import math
import re
import uuid

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Cm,Pt
from matplotlib import pyplot as plt
from pypdf import PdfReader

#调试工具
DEBUG_ON = 1
DEBUG_STAGE = 8

#调试类别
NO_DEBUG_INFO = 1
INFO_DEBUG_INFO = 2
ERROR_DEBUG_INFO = 3
ALL_DEBUG_INFO = 8

#参数定义
GLOBAL_SLIDE_WIDTH=Cm(21)
GLOBAL_SLIDE_HEIGHT=Cm(14.8)

def dprint(DEBUG_LEVEL,*args):
  len_args = range(len(args))
  if DEBUG_ON == 0:
    return
  
  if DEBUG_STAGE == ALL_DEBUG_INFO or DEBUG_STAGE == DEBUG_LEVEL:
    for index in len_args:
      print(args[index])


#预处理
ocr_processor=easyocr.Reader(['ch_sim','en'])

def crop_edge(cvimg,margin):
    print(cvimg.shape)
    return cvimg[margin:cvimg.shape[0]-margin*2,margin:cvimg.shape[1]-margin*2]

def find_split_interval(column_mean):
    """
    找到分割区间。
    column_mean: 每列的平均亮度。
    threshold: 用于确定空白区间的亮度变化阈值。
    """
    data_array = column_mean
    max_value = data_array.max()  # 找到最大值
    max_positions = np.where(data_array == max_value)[0]  # 找到所有最大值的位置

    # 如果最大值只出现一次或两次，那么没有符合条件的连续区间
    if len(max_positions) <= 2:
        return (None, 0)

    # 找到所有最大值连续区间的开始和结束位置
    diff = np.diff(max_positions)
    breaks = np.where(diff > 1)[0]
    starts = np.insert(max_positions[breaks + 1], 0, max_positions[0])
    ends = np.append(max_positions[breaks], max_positions[-1])

    # 排除第一次和最后一次出现的最大值连续区间
    if len(starts) > 2:  # 如果有超过两个连续区间
        inner_lengths = ends[1:-1] - starts[1:-1] + 1
        longest_inner_index = np.argmax(inner_lengths)
        return (starts[1:-1][longest_inner_index], inner_lengths[longest_inner_index])
    else:
        return (None, 0)

def find_optimal_split_position(start, length):
    """
    在确定的区间内找到最合适的分割线。
    """
    # 假定最佳分割点是区间中心
    return start + (length // 2)

def find_split_position(src_img):
    
    gray_image = cv2.cvtColor(src_img, cv2.COLOR_BGR2GRAY)
    
    height, width = gray_image.shape
    
    # 计算每列的像素平均值
    column_sums = np.sum(gray_image, axis=0)
    column_mean = column_sums / height

    if False:
        x = np.linspace(0,width,width)
        y=column_sums
        plt.subplot(2,2,1),
        plt.plot(x,y)
        plt.subplot(2,2,2),
        plt.imshow(src_img,'gray')
        plt.subplot(2,2,3),
        plt.imshow(gray_image,'gray')
        plt.show()
        
    # 找到分割区间
    start, length = find_split_interval(column_sums)
    
    # 计算最优分割位置
    split_position = find_optimal_split_position(start, length)
    
    
    return split_position

def enlarge_foreground_with_grabcut(src_image, scale_factor=3,iterations=5):
    # 读取图像
    image = src_image
    # mask = np.zeros(image.shape[:2], np.uint8)

    original_size = image.shape[:2]

    # 将图像转换为灰度图，并应用一个阈值来区分前景和背景
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    _, mask = cv2.threshold(gray, 254, 255, cv2.THRESH_BINARY_INV)
    
    # 初始化GrabCut算法的掩码和模型参数
    new_mask = np.where(mask == 255, 3, 2).astype('uint8')
    bgdModel = np.zeros((1, 65), np.float64)
    fgdModel = np.zeros((1, 65), np.float64)
    
    # 使用GrabCut算法细化前景和背景分割
    cv2.grabCut(image, new_mask, None, bgdModel, fgdModel, iterations, cv2.GC_INIT_WITH_MASK)
    
    # 将分割结果转换为二值掩码，其中前景为1，背景为0
    mask2 = np.where((new_mask == 2)|(new_mask == 0), 0, 1).astype('uint8')
    
    # 创建一个RGBA通道图像，用于存储带有透明背景的前景
    image_rgba = cv2.cvtColor(image, cv2.COLOR_BGR2BGRA)
    image_rgba[:, :, 3] = mask2 * 255
    
    # 寻找前景的边界，裁剪图像
    x, y, w, h = cv2.boundingRect(mask2)
    cropped_foreground = image_rgba[y:y+h, x:x+w]

    # 放大图像
    enlarged_foreground = cv2.resize(cropped_foreground, None, fx=scale_factor, fy=scale_factor, interpolation=cv2.INTER_CUBIC)
    
    return enlarged_foreground
    # 保存放大后的前景图像

def rearrange_strings(strings):
    # 初始化两个特殊字符串的索引
    english_index = None
    special_char_index = None

    # 遍历字符串数组，寻找满足条件的字符串
    for i, s in enumerate(strings):
        # 找到第一个英文字符串
        if english_index is None and s.isascii():
            english_index = i
        # 找到第一个包含特殊字符的字符串
        if special_char_index is None and ('(' in s or '（' in s):
            special_char_index = i
        # 如果两个条件都已满足，提前结束循环
        if english_index is not None and special_char_index is not None:
            break
    
    # 重新组织字符串数组
    new_strings = []
    # 如果找到英文字符串，添加到新数组第一位
    if english_index is not None:
        new_strings.append(strings.pop(english_index))
    # 如果找到特殊字符字符串，考虑可能已经移除了一个元素的情况
    if special_char_index is not None:
        if english_index is not None and special_char_index > english_index:
            special_char_index -= 1
        new_strings.append(strings.pop(special_char_index))
    # 剩余的字符串按顺序合并后放在第三位
    remaining_string = ''.join(strings)
    if remaining_string:  # 如果还有剩余字符串
        new_strings.append(remaining_string)
    
    return new_strings

def process_slide(slide,img_path,size,textcontent):
    height = size[0]
    width = size[1]
    # new_image_width=Pt(width)
    # new_image_height=Pt(height)
    # new_image_top=(GLOBAL_SLIDE_HEIGHT/2)-(new_image_height/2)
    # new_image_left=(GLOBAL_SLIDE_WIDTH/4)-(new_image_width/2)
    
    print(img_path)
    new_pic_shape = slide.shapes.add_picture(img_path,0,0)

    new_image_width=new_pic_shape.width
    new_image_height=new_pic_shape.height
    new_image_top=int(abs((GLOBAL_SLIDE_HEIGHT/2)-(new_image_height/2)))
    new_image_left=int(abs((GLOBAL_SLIDE_WIDTH/4)-(new_image_width/2)))

    new_pic_shape.top = new_image_top
    new_pic_shape.left = new_image_left

    dprint(INFO_DEBUG_INFO,"****slide size:",GLOBAL_SLIDE_WIDTH/360000,GLOBAL_SLIDE_HEIGHT/360000)
    dprint(INFO_DEBUG_INFO,"****add picture:",new_image_height/360000,new_image_width/360000,new_image_top/360000,new_image_left/360000)
    print("----picture top/left/width/height in number", new_pic_shape.top,new_pic_shape.left,new_pic_shape.width,new_pic_shape.height)
    print("----picture top/left/width/height in CM", new_pic_shape.top/360000,new_pic_shape.left/360000,new_pic_shape.width/360000,new_pic_shape.height/360000)

    txBox = slide.shapes.add_textbox(GLOBAL_SLIDE_WIDTH/2,0,GLOBAL_SLIDE_WIDTH/2,GLOBAL_SLIDE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.clear()

    textcontent = rearrange_strings(textcontent)

    dprint(INFO_DEBUG_INFO,"Adding word to slides")
    p = tf.add_paragraph()
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.text = textcontent[0].lower()
    p.font.name="Calibri"
    p.font.bold=True
    p.font.size=Pt(80)

    p = tf.add_paragraph()
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.text = textcontent[1]
    p.font.name="微软雅黑"
    p.font.bold=False
    p.font.size=Pt(36)

    p = tf.add_paragraph()
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.text = textcontent[2]
    p.font.name="微软雅黑"
    p.font.bold=False
    p.font.size=Pt(36)

       


 
   

   
   
   
def ConvertPdf2PptforEnglishCard(uploaded_files):
  #uploaded_file_name=""
    for uploaded_file in uploaded_files:
        bytes_data = uploaded_file.read()
        uploaded_file_name=uploaded_file.name
        with open(uploaded_file_name, "wb") as f:
            f.write(uploaded_file.getbuffer())

        src_pdf=PdfReader(uploaded_file_name)

        status_message = st.empty()
        progress_bar = st.progress(0)

        status_message.text("开始转换文件："+uploaded_file_name)

        target_pdfile_name = "转换后-" + uploaded_file.name
        target_pptfile_name = os.path.splitext(target_pdfile_name)[0] + '.pptx'
        
        prs_dst = Presentation()
        prs_dst.slide_height=GLOBAL_SLIDE_HEIGHT
        prs_dst.slide_width=GLOBAL_SLIDE_WIDTH
        blank_slide_layout=prs_dst.slide_layouts[6]

        #状态栏标记基数
        pagenum=0
        imgnum=0
        totalImage=len(src_pdf.pages[0].images)

        dprint(INFO_DEBUG_INFO,"总共待处理图片：",totalImage)
        

        #遍历所有页面
        #for eachpdfpage in src_pdf.pages:

        pagenum+=1
        imgnum=0
        #遍历页面中的所有图像元素
        for image_file_object in src_pdf.pages[0].images:

            imgnum+=1
            status_message.text("--开始图像分割与识别，页面:"+str(pagenum)+"图像："+str(imgnum))
            progresspct=int(imgnum/totalImage*100)
            print("当前共有图片：",totalImage,"处理进度到：",progresspct)
            progress_bar.progress(progresspct)

            #将文件临时存储起来
            pdfimg_tmpFile_name="tmpimg"+os.path.splitext(image_file_object.name)[-1]
            with open(pdfimg_tmpFile_name,'wb') as pdf2ppt_tmpFile:
                pdf2ppt_tmpFile.write(image_file_object.data)

            #开始处理图片
            cv2_raw_image = cv2.imread(pdfimg_tmpFile_name)

            #识别图片中的文字
            pdf_text=ocr_processor.readtext(cv2_raw_image,detail=0)

            #识别是卡片图片
            Is_it_a_Card = True
            #for eachword in pdf_text:
            dprint(INFO_DEBUG_INFO,pdf_text[0].find("磨出好耳朵"), pdf_text[0].find("周计划"))

            if pdf_text[0].find("磨出好耳朵")!=-1 or pdf_text[0].find("周计划")!=-1:
                Is_it_a_Card=False
                continue
            dprint(INFO_DEBUG_INFO,pdf_text)
            

            if(Is_it_a_Card==True):
                dprint(INFO_DEBUG_INFO,cv2_raw_image.shape,10)
                cv2_raw_image=crop_edge(cv2_raw_image,10)
                dprint(INFO_DEBUG_INFO,cv2_raw_image.shape)

                split_position = find_split_position(cv2_raw_image)

                dprint(INFO_DEBUG_INFO,split_position)

                # 分割图片
                left_image = cv2_raw_image[:, :split_position]
                #right_image = cv2_raw_image[:, split_position:]

                left_crop_image =enlarge_foreground_with_grabcut(left_image,2,5)
                fileextension=os.path.splitext(image_file_object.name)[1]
                leftimgfilenametmp="Left_img_tmp"+fileextension
                cv2.imwrite(leftimgfilenametmp,left_crop_image)

                slide_dst = prs_dst.slides.add_slide(blank_slide_layout)
                print(left_crop_image.shape)
                process_slide(slide_dst,leftimgfilenametmp,left_crop_image.shape,pdf_text)

        prs_dst.save(target_pptfile_name)

        progress_bar.progress(100)
        status_message.text("转换完成！"+uploaded_file_name)
        st.session_state.download_files.append(target_pptfile_name)
        #prs_dst.saveas(target_pdffile_name,FileFormat=32)
        # with open(target_pptfile_name, 'rb') as pptf:
        #     target_file = pptf
        #     if st.download_button('下载转换后的pptx', target_file.read(),file_name=target_pptfile_name,mime="pptx",key=uuid.uuid1()):
        #         status_message.text("开始下载！"+uploaded_file_name)
    st.session_state.uploaded_files = uploaded_files
      # with open(target_pdffile_name, 'rb') as pdff:
      #   target_file = pdff
      # st.download_button('下载转换后的pdf', target_file.read(),file_name=target_pdffile_name,mime="pdf") 




def show_web_icon():
    st.image("https://github.com/dragonleehom/pdf2/raw/master/data/icon.png",width=200)
  

# Page title
st.set_page_config(page_title='<center>铭铭的英语学习卡转换器<center>', page_icon='https://github.com/dragonleehom/pdf2/raw/master/data/icon.png')
#st.image('data/Ming.png', caption='')
# 创建三列，图像放在中间列

show_web_icon()
st.title('        铭铭的英语学习卡转换器        ')
#st.markdown("<center>铭铭的英语学习卡转换器<center>", unsafe_allow_html=True)

with st.expander('关于工具'):
  st.markdown('**工具能做什么?**')
  st.info('本工具是将PDF形式的英语学习卡片进行拆分，并生成一个PPT格式的文档用于打印成卡片')
  st.markdown('**如何使用?**')
  st.warning('请将特定的PDF文本直接拖入到文件上传框中即可.')
  
st.subheader('请上传待转换文件,仅支持pdf,可同时上传多个文件')

# 初始化会话状态
if 'computation_done' not in st.session_state:
    st.session_state.computation_done = False
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = []
if 'repeatrun' not in st.session_state:
    st.session_state.repeatrun = 0
if 'download_files' not in st.session_state:
    st.session_state.download_files = []

st.session_state.repeatrun += 1

uploaded_files=st.file_uploader("请上传文件",type=['pdf','PDF'],accept_multiple_files=True)


if uploaded_files is not None and uploaded_files != st.session_state.uploaded_files:
  st.session_state.computation_done = False
  #with st.status("开始启动转换...", expanded=True) as status:
  st.session_state.uploaded_files = uploaded_files
  ConvertPdf2PptforEnglishCard(uploaded_files)

for each_download_file in st.session_state.download_files:
    with open(each_download_file, 'rb') as pptf:
        target_file = pptf
        st.write(each_download_file) 
        st.download_button('下载转换后的pptx', target_file.read(),file_name=each_download_file,mime="pptx",key=uuid.uuid1())



